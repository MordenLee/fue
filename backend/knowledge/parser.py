"""File parsing utilities — extract plain text from a wide range of file formats.

Supported formats
-----------------
Plain text  : .txt  .md  .csv  .json
PDF         : .pdf   — engines: pdfplumber (default) | pymupdf | pypdf
Word        : .docx  — engines: python-docx (default) | markitdown
              .doc   — markitdown (tries python-docx as fallback)
PowerPoint  : .pptx
Excel       : .xlsx  .xls
Web / markup: .html  .htm  .xml
E-book      : .epub
Rich text   : .rtf
OpenDocument: .odt
"""

from __future__ import annotations

from pathlib import Path

# ---------------------------------------------------------------------------
# Public interface
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    ".txt", ".md",
    ".pdf",
    ".docx", ".doc",
    ".pptx",
    ".xlsx", ".xls",
    ".csv",
    ".html", ".htm",
    ".xml",
    ".epub",
    ".rtf",
    ".odt",
    ".json",
}


def parse_file(
    file_path: str,
    pdf_parser: str = "pdfplumber",
    docx_parser: str = "python-docx",
) -> str:
    """Read *file_path* and return its text content.

    Parameters
    ----------
    file_path:   Absolute path to the file on disk.
    pdf_parser:  "pdfplumber" | "pymupdf" | "pypdf"
    docx_parser: "python-docx" | "markitdown"
    """
    p = Path(file_path)
    suffix = p.suffix.lower()

    if suffix in (".txt", ".md"):
        return p.read_text(encoding="utf-8", errors="replace")
    if suffix == ".csv":
        return _parse_csv(p)
    if suffix == ".json":
        return _parse_json(p)
    if suffix in (".html", ".htm"):
        return _parse_html(p)
    if suffix == ".xml":
        return _parse_xml(p)
    if suffix == ".pdf":
        return _parse_pdf(p, pdf_parser)
    if suffix == ".docx":
        return _parse_docx(p, docx_parser)
    if suffix == ".doc":
        return _parse_doc(p)
    if suffix == ".pptx":
        return _parse_pptx(p)
    if suffix == ".xlsx":
        return _parse_xlsx(p)
    if suffix == ".xls":
        return _parse_xls(p)
    if suffix == ".epub":
        return _parse_epub(p)
    if suffix == ".rtf":
        return _parse_rtf(p)
    if suffix == ".odt":
        return _parse_odt(p)

    raise ValueError(f"Unsupported file type: {suffix}")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _parse_pdf(path: Path, engine: str) -> str:
    if engine == "pymupdf":
        return _pdf_pymupdf(path)
    if engine == "pypdf":
        return _pdf_pypdf(path)
    return _pdf_pdfplumber(path)  # default: pdfplumber


def _pdf_pdfplumber(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    return "\n".join(parts)


def _pdf_pymupdf(path: Path) -> str:
    import fitz  # PyMuPDF

    parts: list[str] = []
    with fitz.open(str(path)) as doc:
        for page in doc:
            parts.append(page.get_text())
    return "\n".join(parts)


def _pdf_pypdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(p for p in parts if p.strip())


# ---------------------------------------------------------------------------
# Word — DOCX
# ---------------------------------------------------------------------------

def _parse_docx(path: Path, engine: str) -> str:
    if engine == "markitdown":
        return _markitdown_convert(path)
    return _docx_python_docx(path)  # default


def _docx_python_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Word — legacy binary .doc
# ---------------------------------------------------------------------------

def _parse_doc(path: Path) -> str:
    """Legacy binary .doc — try markitdown first, then python-docx as fallback."""
    try:
        return _markitdown_convert(path)
    except Exception:
        pass
    try:
        return _docx_python_docx(path)
    except Exception as exc:
        raise ValueError(
            f"Cannot parse '{path.name}'. "
            "Old binary .doc format has limited pure-Python support — "
            "please convert to .docx for best results."
        ) from exc


def _markitdown_convert(path: Path) -> str:
    from markitdown import MarkItDown

    md = MarkItDown()
    result = md.convert(str(path))
    return result.text_content or ""


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def _parse_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def _parse_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(c) for c in row if c is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _parse_xls(path: Path) -> str:
    import xlrd

    wb = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for ws in wb.sheets():
        for i in range(ws.nrows):
            row_text = "\t".join(str(ws.cell_value(i, j)) for j in range(ws.ncols))
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# HTML / XML
# ---------------------------------------------------------------------------

def _parse_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_bytes(), "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_xml(path: Path) -> str:
    import xml.etree.ElementTree as ET

    tree = ET.parse(str(path))
    root = tree.getroot()
    texts = [t.strip() for t in root.itertext() if t.strip()]
    return "\n".join(texts)


# ---------------------------------------------------------------------------
# EPUB
# ---------------------------------------------------------------------------

def _parse_epub(path: Path) -> str:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup

    book = epub.read_epub(str(path), options={"ignore_ncx": True})
    parts: list[str] = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "lxml")
        text = soup.get_text(separator="\n", strip=True)
        if text:
            parts.append(text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# RTF
# ---------------------------------------------------------------------------

def _parse_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text

    content = path.read_text(encoding="utf-8", errors="replace")
    return rtf_to_text(content)


# ---------------------------------------------------------------------------
# ODT (OpenDocument Text)
# ---------------------------------------------------------------------------

def _parse_odt(path: Path) -> str:
    from odf.opendocument import load as odf_load
    from odf.text import P

    doc = odf_load(str(path))
    parts: list[str] = []

    def _extract(node) -> str:
        if node.nodeType == node.TEXT_NODE:
            return node.data
        return "".join(_extract(child) for child in node.childNodes)

    for para in doc.getElementsByType(P):
        text = _extract(para).strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# CSV / JSON
# ---------------------------------------------------------------------------

def _parse_csv(path: Path) -> str:
    import csv
    import io

    content = path.read_text(encoding="utf-8-sig", errors="replace")
    reader = csv.reader(io.StringIO(content))
    rows = ["\t".join(cell.strip() for cell in row) for row in reader if any(row)]
    return "\n".join(rows)


def _parse_json(path: Path) -> str:
    import json

    data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    return json.dumps(data, ensure_ascii=False, indent=2)

